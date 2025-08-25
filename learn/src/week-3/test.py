from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.shapes import MSO_CONNECTOR
from pptx.dml.color import RGBColor
from pptx.enum.dml import MSO_LINE_DASH_STYLE

# ===== بيانات الرسم =====
PROBLEM_TEXT = "المشكلة:\nعدم تنظيف وتعقيم الأسطح"
FACTORS = {
    "العوامل البشرية (Man)": [
        "ضعف الوعي بأهمية التعقيم",
        "قلة المتابعة والإشراف",
        "تقصير في الالتزام بالتعليمات",
    ],
    "الطرق (Methods)": [
        "عدم وجود بروتوكول ثابت للتنظيف",
        "ضعف نظام المراقبة",
        "عدم وجود خطة دورية للتنظيف",
    ],
    "المواد (Materials)": [
        "نقص في المواد المطهرة",
        "استخدام مواد غير فعالة",
        "انتهاء صلاحية بعض المواد",
    ],
    "البيئة (Environment)": [
        "ازدحام داخل الوحدة",
        "ضغط العمل وكثرة المرضى",
        "بيئة غير مهيأة للعمل",
    ],
    "الأجهزة (Equipment)": [
        "تعطل بعض الأجهزة",
        "عدم توفر أدوات تعقيم كافية",
        "ضعف الصيانة الدورية",
    ],
}

# ===== إنشاء الشريحة =====
prs = Presentation()
slide = prs.slides.add_slide(prs.slide_layouts[6])  # شريحة فارغة

# عنوان
title = slide.shapes.add_textbox(Inches(1.2), Inches(0.3), Inches(7.5), Inches(0.8))
t = title.text_frame
t.text = "مخطط عظمة السمكة (Fishbone / Ishikawa)"
t.paragraphs[0].font.size = Pt(26)

# العمود الفقري (سهم طويل من اليمين لليسار)
spine = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(0.7), Inches(3.1), Inches(8.3), Inches(0.6))
spine.fill.solid();
spine.fill.fore_color.rgb = RGBColor(230, 230, 230)
spine.line.color.rgb = RGBColor(50, 50, 50)

# نص المشكلة على رأس السهم (يسار)
problem = slide.shapes.add_textbox(Inches(7.9), Inches(2.8), Inches(2.1), Inches(1.1))
tf = problem.text_frame
tf.word_wrap = True
tf.text = PROBLEM_TEXT
tf.paragraphs[0].font.size = Pt(20)
for r in tf.paragraphs[0].runs: r.font.bold = True

TITLE_COLOR = RGBColor(79, 129, 189)
BULLET_SIZE = Pt(14)
TITLE_SIZE = Pt(16)

# أماكن صناديق الفروع (أعلى/أسفل ويمين/يسار)
# (x, y, w, h)
boxes = {
    "العوامل البشرية (Man)": (Inches(0.2), Inches(1.1), Inches(3.2), Inches(1.6)),
    "المواد (Materials)": (Inches(3.9), Inches(1.1), Inches(3.2), Inches(1.6)),
    "الطرق (Methods)": (Inches(0.2), Inches(4.6), Inches(3.2), Inches(1.6)),
    "البيئة (Environment)": (Inches(3.9), Inches(4.6), Inches(3.2), Inches(1.6)),
    "الأجهزة (Equipment)": (Inches(2.9), Inches(5.9), Inches(3.2), Inches(1.6)),
}

# نقاط اتصال على العمود الفقري (لعمل عظام مائلة)
spine_points = {
    "العوامل البشرية (Man)": (Inches(2.1), Inches(3.1)),  # فوق
    "المواد (Materials)": (Inches(4.8), Inches(3.1)),  # فوق
    "الطرق (Methods)": (Inches(2.1), Inches(3.7)),  # تحت
    "البيئة (Environment)": (Inches(6.2), Inches(3.7)),  # تحت
    "الأجهزة (Equipment)": (Inches(5.6), Inches(3.4)),  # وسط
}

# رسم صناديق الفروع وتعبئتها
shape_boxes = {}
for title, (x, y, w, h) in boxes.items():
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    box.fill.solid();
    box.fill.fore_color.rgb = RGBColor(242, 246, 252)
    box.line.color.rgb = RGBColor(100, 100, 100)
    tf = box.text_frame;
    tf.clear()

    # العنوان
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = TITLE_SIZE
    p.font.bold = True
    p.font.color.rgb = TITLE_COLOR

    # النقاط
    for item in FACTORS[title]:
        p = tf.add_paragraph()
        p.text = "• " + item
        p.level = 1
        p.font.size = BULLET_SIZE

    shape_boxes[title] = box


# دالة رسم خط/موصل بين نقطتين
def add_connector(x1, y1, x2, y2, dashed=False):
    conn = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, x1, y1, x2, y2)
    conn.line.color.rgb = RGBColor(70, 70, 70)
    conn.line.width = Pt(2)
    if dashed:
        conn.line.dash_style = MSO_LINE_DASH_STYLE.DASH
    return conn


# توصيل الصناديق بالعمود الفقري بخطوط مائلة
for title, (sx, sy) in spine_points.items():
    box = shape_boxes[title]
    # لو الصندوق على يسار العمود نخرج من يمينه، والعكس
    if box.left < Inches(3.5):
        bx = box.left + box.width
    else:
        bx = box.left
    by = box.top + box.height / 2
    add_connector(bx, by, sx, sy)

# حفظ الملف
out_path = "مخطط_عظمة_السمكة.pptx"
prs.save(out_path)
print(f"تم إنشاء الملف: {out_path}")